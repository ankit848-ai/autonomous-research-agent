import os
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase.pdfmetrics import stringWidth
from pptx import Presentation
from pptx.util import Inches

class Writer:

    def create_outputs(self, literature, plan, analysis, session_id):
        folder = f"artifacts/{session_id}"
        os.makedirs(folder, exist_ok=True)

        # -------------------- CREATE PDF --------------------
        pdf_path = f"{folder}/paper.pdf"
        styles = getSampleStyleSheet()
        story = []

        doc = SimpleDocTemplate(pdf_path, pagesize=A4)

        # Title
        story.append(Paragraph("<b>Research Paper</b>", styles["Title"]))
        story.append(Spacer(1, 12))

        # Experiment info
        story.append(Paragraph(f"<b>Topic:</b> {plan['topic']}", styles["BodyText"]))
        story.append(Paragraph(f"<b>Model:</b> {plan['model']}", styles["BodyText"]))
        story.append(Paragraph(f"<b>Metric:</b> {plan['metric']}", styles["BodyText"]))
        story.append(Paragraph(f"<b>Accuracy:</b> {analysis['accuracy']}", styles["BodyText"]))
        story.append(Spacer(1, 12))

        # Literature
        story.append(Paragraph("<b>Literature Summary:</b>", styles["Heading2"]))
        for p in literature:
            story.append(Paragraph(f"- <b>{p['title']}</b>: {p['summary']}", styles["BodyText"]))
            story.append(Spacer(1, 6))

        doc.build(story)

        # -------------------- CREATE PPT --------------------
        ppt_path = f"{folder}/slides.pptx"
        prs = Presentation()

        # Slide 1
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "ResearchLab-AI Results"

        subtitle = slide.placeholders[1]
        subtitle.text = "Automatically generated research presentation"

        # Slide 2 - Experiment Overview
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Experiment Overview"
        body = slide2.placeholders[1]
        body.text = (
            f"Topic: {plan['topic']}\n"
            f"Model: {plan['model']}\n"
            f"Metric: {plan['metric']}\n"
            f"Accuracy: {analysis['accuracy']}"
        )

        # Slide 3 - Literature
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])
        slide3.shapes.title.text = "Literature Summary"
        slide3.placeholders[1].text = "\n".join([f"- {p['title']}" for p in literature])

        # Slide 4 - Accuracy Plot
        slide4 = prs.slides.add_slide(prs.slide_layouts[1])
        slide4.shapes.title.text = "Accuracy Plot"
        pic = slide4.shapes.add_picture(analysis["plot_path"], Inches(1), Inches(1), width=Inches(5))

        prs.save(ppt_path)

        return {
            "paper": pdf_path,
            "slides": ppt_path
        }
